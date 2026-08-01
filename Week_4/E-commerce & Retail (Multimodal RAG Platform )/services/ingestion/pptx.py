from pptx import Presentation


def extract_pptx(file_path: str) -> list[str]:
    """
    Extract text from all slides in a PowerPoint presentation.
    Returns:
        list[str]: One string per slide.
    """

    presentation = Presentation(file_path)
    texts = []

    for slide in presentation.slides:
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        if slide_text:
            texts.append("\n".join(slide_text))

    return texts